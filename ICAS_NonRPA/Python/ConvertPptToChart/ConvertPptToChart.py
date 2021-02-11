'''
Copyright 2020 Infosys Ltd.
Use of this source code is governed by Apache 2.0 license that can be found in the LICENSE file or at 
https://opensource.org/licenses/Apache-2.0 .
'''
from pptx import Presentation
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
from abstract_bot import Bot

# Python Bot to convert tables of PPT document into pie chart

class ConvertPptToChart(Bot):
	
	def divide_chunks(self, l, n): 
		for i in range(0, len(l), n):       # divide the list into equal parts
			yield l[i:i + n]

	def execute(self, executeContext) :
		try:
			pptFilePath = executeContext['pptFilePath']
			if  not pptFilePath:
				return {'validation error' : 'missing argument pptFilePath'}

			outputDestination = executeContext['outputDestination']
			if  not outputDestination:
				return {'validation error' : 'missing argument outputDestination'}

			outputChartFileName = executeContext['outputChartFileName']
			if  not outputChartFileName:
				return {'validation error' : 'missing argument outputChartFileName'}

			destination_path = '{0}\{1}'.format(outputDestination, outputChartFileName)
	
			pptFile = Presentation(pptFilePath)
			chartPDF = PdfPages(destination_path)

			main_list =[]
			columns_count = []
			
			for slide in pptFile.slides:
				mini_list = []
				for shape in slide.shapes:
					if not shape.has_table:
						continue    
					tbl = shape.table
					row_count = len(tbl.rows)
					col_count = len(tbl.columns)
					columns_count.append(col_count)
					for r in range(0, row_count):
						for c in range(0, col_count):
							cell = tbl.cell(r,c)
							paragraphs = cell.text_frame.paragraphs 
							for paragraph in paragraphs:
								for run in paragraph.runs:
									digit_check  = run.text
									value = run.text
									if digit_check.isdigit():
										mini_list.append(int(value))
									elif digit_check.replace('.', '', 1).isdigit():
										mini_list.append(float(value))
									else:
										mini_list.append(value)
				if mini_list:
					main_list.append(mini_list)
			#print(main_list)

			##############################
			# dividing the list with 
			# headers and column totals
			##############################
			new_data_set = []
			counter = 0
			for item in main_list:
				data = list(self.divide_chunks(item, columns_count[counter]))
				counter = counter + 1
				new_data_set.append(data)
			#print (new_data_set)

			##################################
			# creating headers
			##################################
			headers = [heads[0] for heads  in new_data_set]
			#print (headers)

			##################################
			# creating a new list by removing 
			# string headers for addition of   
			# the list of arrays for pie chart
			##################################
			for data in new_data_set:
				data.pop(0)

			##################################
			# adding the list of array
			##################################
			calculated_data = []
			for array_list in new_data_set:
				res = [sum(i) for i in zip(*array_list)] 
				calculated_data.append(res)
			#print (calculated_data)


			# ##### plotting pie-chart #############
			fig = plt.figure(3, figsize=(4,4))
			ax = plt.axes([0.1, 0.1, 0.8, 0.8])
			for item in range (0, len(calculated_data)):
				ax.cla()    
				labels=headers[item]       
				fracs=calculated_data[item]   
				p = plt.pie(fracs,labels=labels,autopct='%1.1f%%')
				fig.savefig(chartPDF, format='pdf')   

			chartPDF.close()
			return {'status' : 'success'}
		except Exception as e:
			#print("Error occured",str(e)) 
			return {'Exception' : str(e)}

  
if __name__ == '__main__':
	context = {}
	bot_obj = ConvertPptToChart()

	context = {
				
				'pptFilePath': "",
				'outputDestination': "", 
				'outputChartFileName' : "" 

			}
	bot_obj.bot_init()
	output = bot_obj.execute(context)
	print('output : ',output)







